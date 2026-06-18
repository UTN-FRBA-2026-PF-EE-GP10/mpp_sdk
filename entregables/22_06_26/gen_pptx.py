#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.10"
# dependencies = [
#     "python-pptx>=0.6.23",
# ]
# ///
"""Genera la presentacion del anteproyecto (22/06/26) en formato .pptx.

El archivo de salida se puede subir directamente a Google Slides.

Uso:
    uv run entregables/22_06_26/gen_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
IMG = HERE / "img"

# Paleta UTN / Google-kanban
AZUL = RGBColor(0x1A, 0x23, 0x7E)
GRIS = RGBColor(0x55, 0x55, 0x55)
VERDE = RGBColor(0x34, 0xA8, 0x53)
NEGRO = RGBColor(0x20, 0x20, 0x20)

EMU16x9_W = Inches(13.333)
EMU16x9_H = Inches(7.5)


def _slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def _box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _set(p, text, size, color=NEGRO, bold=False, align=PP_ALIGN.LEFT, italic=False):
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def title_slide(prs):
    s = _slide(prs)
    tf = _box(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(2.0))
    _set(
        tf.paragraphs[0],
        "Framework para Comparacion y Desarrollo de Algoritmos MPPT",
        34,
        AZUL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    p = tf.add_paragraph()
    _set(
        p,
        "Tesis Final de Grado - Ingenieria Electronica",
        18,
        GRIS,
        align=PP_ALIGN.CENTER,
        italic=True,
    )
    tf2 = _box(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(1.5))
    _set(
        tf2.paragraphs[0],
        "Federico Borello   |   Juan Falabella   |   Diego Nahuel Pirotta",
        16,
        NEGRO,
        align=PP_ALIGN.CENTER,
    )
    p = tf2.add_paragraph()
    _set(p, "UTN - Facultad Regional Buenos Aires", 14, GRIS, align=PP_ALIGN.CENTER)
    p = tf2.add_paragraph()
    _set(p, "22 de junio de 2026", 14, GRIS, align=PP_ALIGN.CENTER)


def bullets_slide(prs, title, bullets, subtitle=None):
    s = _slide(prs)
    tf = _box(s, Inches(0.7), Inches(0.4), Inches(12), Inches(1.0))
    _set(tf.paragraphs[0], title, 30, AZUL, bold=True)
    if subtitle:
        p = tf.add_paragraph()
        _set(p, subtitle, 16, GRIS, italic=True)
    body = _box(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.3))
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        _set(p, "-  " + b, 20, NEGRO)
        p.space_after = Pt(12)


def image_slide(prs, title, images, caption=None, note=None):
    """images: lista de (path, left, top, width)."""
    s = _slide(prs)
    tf = _box(s, Inches(0.7), Inches(0.35), Inches(12), Inches(0.9))
    _set(tf.paragraphs[0], title, 28, AZUL, bold=True)
    if note:
        nf = _box(s, Inches(0.7), Inches(1.15), Inches(12), Inches(0.5))
        _set(nf.paragraphs[0], note, 14, VERDE, italic=True)
    for path, left, top, width in images:
        s.shapes.add_picture(str(path), left, top, width=width)
    if caption:
        cf = _box(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.5))
        _set(cf.paragraphs[0], caption, 13, GRIS, italic=True, align=PP_ALIGN.CENTER)


def table_slide(prs, title, headers, rows, note=None, caption=None):
    """Tabla simple. headers: lista de str; rows: lista de listas de str."""
    s = _slide(prs)
    tf = _box(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
    _set(tf.paragraphs[0], title, 28, AZUL, bold=True)
    if note:
        nf = _box(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.5))
        _set(nf.paragraphs[0], note, 16, VERDE, italic=True)
    nrows, ncols = len(rows) + 1, len(headers)
    left, top = Inches(0.9), Inches(2.0)
    width, height = Inches(11.5), Inches(0.5 * nrows)
    tbl = s.shapes.add_table(nrows, ncols, left, top, width, height).table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = h
        para = c.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        r = para.runs[0]
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        c.fill.solid()
        c.fill.fore_color.rgb = AZUL
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.text = str(val)
            para = c.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = para.runs[0]
            r.font.size = Pt(15)
            r.font.color.rgb = NEGRO
    if caption:
        cf = _box(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.5))
        _set(cf.paragraphs[0], caption, 13, GRIS, italic=True, align=PP_ALIGN.CENTER)


def split_slide(prs, title, bullets, image, img_left, img_top, img_width):
    s = _slide(prs)
    tf = _box(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.9))
    _set(tf.paragraphs[0], title, 28, AZUL, bold=True)
    body = _box(s, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.2))
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        _set(p, "-  " + b, 18, NEGRO)
        p.space_after = Pt(10)
    s.shapes.add_picture(str(image), img_left, img_top, width=img_width)


def main() -> None:
    prs = Presentation()
    prs.slide_width = EMU16x9_W
    prs.slide_height = EMU16x9_H

    title_slide(prs)

    bullets_slide(
        prs,
        "El problema",
        [
            "Los paneles FV necesitan algoritmos MPPT para extraer la maxima potencia bajo irradiancia y temperatura variables.",
            "Existen muchos metodos (P&O, InCond, Fuzzy, PSO...) pero se comparan en MATLAB con codigo rara vez publicado.",
            "Pasar de la simulacion al hardware suele exigir reescribir la logica de control, introduciendo fricciones y discrepancias.",
            "Falta una biblioteca de control abierta con un banco de comparacion uniforme.",
        ],
    )

    bullets_slide(
        prs,
        "Enfoque propuesto",
        [
            "Una capa de abstraccion de senales separa el algoritmo del medio de ejecucion.",
            "El algoritmo solo ve: read() -> (V, I) y write(D).",
            "El mismo codigo corre sobre un panel simulado (pvlib) o sobre el SEPIC real, cambiando solo el adaptador.",
            "Esa frontera es tambien la del firmware: el RP2040 es otra fuente de senales para la Raspberry Pi.",
            "Convertidor SEPIC: eleva o reduce la tension del panel sin cambiar de topologia.",
        ],
        subtitle="Simulacion y hardware con el mismo codigo y las mismas metricas",
    )

    split_slide(
        prs,
        "SDK de simulacion y metricas",
        [
            "5 algoritmos bajo la misma interfaz: P&O, InCond, Fuzzy, Scan-and-Track, PSO.",
            "Locales vs globales: los globales escapan de los maximos locales del sombreado parcial.",
            "Cada algoritmo es portable a microcontrolador (sin dependencias, poco estado).",
            "Metrica principal: eficiencia energetica en el tiempo.",
        ],
        IMG / "sim_cyclic.png",
        Inches(6.7),
        Inches(1.5),
        Inches(6.2),
    )

    image_slide(
        prs,
        "Herramienta interactiva de comparacion en vivo",
        [
            (IMG / "animate.png", Inches(1.4), Inches(1.5), Inches(10.5)),
        ],
        caption="Sombreado parcial en tiempo real (sliders de irradiancia): la curva P-V tiene dos picos; un metodo queda anclado al pico local (~88%) mientras el resto alcanza el MPP global.",
    )

    split_slide(
        prs,
        "Robustez al ruido de medicion",
        [
            "El ruido se modela como una capa que envuelve a la fuente y solo perturba lo que el controlador ve.",
            "Los rastreadores locales de paso fijo colapsan cuando el ruido supera el cambio de potencia por perturbacion.",
            "Los metodos globales sostienen un piso gracias a sus busquedas periodicas.",
            "Motiva un paso adaptativo y un disparo de re-busqueda en el algoritmo propio.",
        ],
        IMG / "sim_noise.png",
        Inches(6.7),
        Inches(1.6),
        Inches(6.2),
    )

    table_slide(
        prs,
        "Metricas - perfil ciclico (seed 1)",
        ["Algoritmo", "eta", "Acierto", "Atrap.", "P.atrap.", "reacq"],
        [
            ["P&O", "88,2%", "16/30", "14", "65,4%", "9 ms"],
            ["InCond", "88,3%", "16/30", "14", "65,4%", "8 ms"],
            ["Fuzzy", "88,5%", "16/30", "14", "65,6%", "7 ms"],
            ["Scan&Track", "95,0%", "25/30", "5", "86,8%", "464 ms"],
            ["PSO", "93,4%", "21/30", "9", "88,0%", "483 ms"],
        ],
        note="Metricas iniciales - seguimos trabajando en esto.",
        caption="Locales: rapidos pero atrapados en casi la mitad de los plateaus. Globales: casi no se trampean, a costa de re-adquisicion mas lenta.",
    )

    table_slide(
        prs,
        "Metricas - robustez al ruido (eta energia)",
        ["sigma %FE", "P&O", "InCond", "Fuzzy", "S&T", "PSO"],
        [
            ["0,00", "91,0%", "91,0%", "91,3%", "94,6%", "94,2%"],
            ["0,25", "87,6%", "87,6%", "90,0%", "92,1%", "91,5%"],
            ["0,50", "55,6%", "55,5%", "82,8%", "80,8%", "79,7%"],
            ["1,00", "19,3%", "19,5%", "46,2%", "64,1%", "60,7%"],
            ["2,00", "15,2%", "15,4%", "18,8%", "62,1%", "60,5%"],
        ],
        note="Metricas iniciales - seguimos trabajando en esto.",
        caption="Los locales de paso fijo colapsan entre 0,25 y 0,50 %FE; los globales sostienen un piso cercano al 60%.",
    )

    image_slide(
        prs,
        "Placa de potencia",
        [
            (IMG / "pcb_real.png", Inches(0.7), Inches(1.5), Inches(6.0)),
            (IMG / "pcb_3d.png", Inches(7.0), Inches(1.6), Inches(5.7)),
        ],
        caption="PCB v1.0 fabricado (SEPIC + driver de gate + INA226 + RP2040). Componentes comprados, armado en curso.",
    )

    image_slide(
        prs,
        "Banco de trabajo preliminar",
        [
            (IMG / "setup_preliminar.png", Inches(0.7), Inches(1.7), Inches(6.0)),
            (IMG / "firmware_setup.png", Inches(7.0), Inches(1.7), Inches(5.7)),
        ],
        caption="Forma preliminar de probar el setup, principalmente los paneles (2x Hissuma PSF10MONO en serie). Comunicacion SPI RP2040 <-> Raspberry Pi validada.",
    )

    image_slide(
        prs,
        "Simulacion conmutada en PLECS",
        [
            (IMG / "plecs.png", Inches(2.4), Inches(1.6), Inches(8.5)),
        ],
        caption="PLECS valida el hardware (ciclo a ciclo); el SDK valida el algoritmo. Definir como comparar ambos es trabajo en curso.",
    )

    bullets_slide(
        prs,
        "Algoritmo propio",
        [
            "El aporte central es el framework sim-a-real, no un algoritmo nuevo: el espacio MPPT ya esta muy explorado.",
            "El algoritmo propio es un caso de estudio para poner a prueba el framework.",
            "Linea en evaluacion: Scan-and-Track + disparo de re-busqueda informado (detector |dP|/P + re-busqueda periodica).",
            "Etapa abierta y exploratoria: probablemente no se implemente algo nuevo, sino experimentar variantes de bajo costo.",
            "La pregunta: alcanza un metodo barato (apto MCU) la eficiencia de los globales pesados? El resultado sirve gane o pierda.",
            "Todo sera un compromiso entre eficiencia de rastreo y costo de implementacion.",
        ],
    )

    bullets_slide(
        prs,
        "Estado y plan de trabajo",
        [
            "Listo: PCB fabricado, componentes comprados, firmware SPI inicial, banco preliminar, SDK con 5 algoritmos y metricas.",
            "En progreso: armado de la placa (~2 semanas), cierre de simulaciones, PLECS, comparacion PLECS vs SDK.",
            "Por hacer: verificacion de la placa, firmware de control (fin de junio / mediados de julio), cierre del algoritmo propio, validacion sim-a-real.",
        ],
    )

    s = _slide(prs)
    tf = _box(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.5))
    _set(tf.paragraphs[0], "Gracias", 40, AZUL, bold=True, align=PP_ALIGN.CENTER)
    p = tf.add_paragraph()
    _set(p, "Proyecto MPPT - UTN FRBA", 18, GRIS, align=PP_ALIGN.CENTER)

    out = HERE / "presentacion_22_06_26.pptx"
    prs.save(str(out))
    print(f"Guardado: {out}")


if __name__ == "__main__":
    main()
